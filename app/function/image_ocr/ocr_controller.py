import os
import json
import string
import tempfile
import shutil
from typing import List, Dict, Optional, Tuple, Any
from pathlib import Path
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import os
import re
sys.path.insert(0, os.path.dirname(__file__))

# 添加QwenTranslator导入
from .translator import QwenTranslator

def perform_ocr_on_image(image_path: str, api_key: str) -> Optional[Dict]:
    """
    对单个图片执行OCR识别
    
    Args:
        image_path: 图片路径
        api_key: OCR API密钥
        
    Returns:
        OCR结果字典或None
    """
    logger.info(f"[OCR识别] 开始对图片执行OCR识别: {image_path}")
    logger.info(f"[OCR识别] 图片文件大小: {os.path.getsize(image_path) if os.path.exists(image_path) else '文件不存在'} 字节")
    
    try:
        # 检查文件是否存在
        if not os.path.exists(image_path):
            logger.error(f"[OCR识别] 图片文件不存在: {image_path}")
            return None
            
        # 检查API密钥
        if not api_key:
            logger.error("[OCR识别] API密钥未提供")
            return None
            
        # 读取图片文件
        with open(image_path, 'rb') as f:
            image_data = f.read()
            
        logger.info(f"[OCR识别] 成功读取图片数据，大小: {len(image_data)} 字节")
        
        # 调用OCR API
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/octet-stream'
        }
        
        # 根据文件扩展名确定OCR类型
        _, ext = os.path.splitext(image_path)
        ext = ext.lower()
        
        if ext in ['.pdf']:
            ocr_type = 'pdf'
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif']:
            ocr_type = 'image'
        else:
            ocr_type = 'image'  # 默认为图片
            
        logger.info(f"[OCR识别] 图片类型: {ext}, OCR类型: {ocr_type}")
        
        # 构造API URL
        api_url = "https://dashscope.aliyuncs.com/api/v1/services/ocr/general"
        params = {'ocr_type': ocr_type}
        
        logger.info(f"[OCR识别] 调用OCR API: {api_url}")
        logger.info(f"[OCR识别] 请求参数: {params}")
        
        # 发送请求
        response = requests.post(
            api_url,
            headers=headers,
            params=params,
            data=image_data,
            timeout=300  # 5分钟超时
        )
        
        logger.info(f"[OCR识别] OCR API响应状态码: {response.status_code}")
        
        # 检查响应
        if response.status_code == 200:
            result = response.json()
            logger.info(f"[OCR识别] OCR API响应成功")
            logger.info(f"[OCR识别] OCR结果结构: {list(result.keys()) if isinstance(result, dict) else '非字典结构'}")
            
            # 提取文本信息
            if 'output' in result and 'recognized_texts' in result['output']:
                texts = result['output']['recognized_texts']
                logger.info(f"[OCR识别] 识别到 {len(texts)} 个文本块")
                
                # 构造返回结果
                extracted_texts = {}
                for i, text_info in enumerate(texts):
                    if isinstance(text_info, dict) and 'text' in text_info:
                        text_key = f"text_{i+1:03d}"
                        extracted_texts[text_key] = text_info['text']
                        logger.info(f"[OCR识别] 文本 {text_key}: {text_info['text'][:50]}...")
                
                return {
                    'texts': extracted_texts,
                    'raw_result': result
                }
            else:
                logger.warning(f"[OCR识别] OCR结果格式不符合预期: {list(result.keys()) if isinstance(result, dict) else '非字典结构'}")
                return {
                    'texts': {},
                    'raw_result': result
                }
        else:
            logger.error(f"[OCR识别] OCR API调用失败，状态码: {response.status_code}")
            logger.error(f"[OCR识别] 错误响应: {response.text}")
            return None
            
    except Exception as e:
        logger.error(f"[OCR识别] OCR识别过程中出错: {str(e)}")
        logger.exception("[OCR识别] 详细错误信息")
        return None
# 导入日志系统
from logger_config_ocr import get_logger

# 导入OCR API处理器
from ocr_api import OCRProcessor

# 导入OCR QWEN API处理程序
from qwen_ocr_api import process_folder_with_mapping

# 导入翻译模块
from translator import TranslationManager

# 获取日志记录器
logger = get_logger("ocr_controller")


def _extract_image_paths_from_markdown(markdown_content: str, markdown_dir: str) -> List[str]:
    """
    从Markdown内容中提取图片路径
    
    Args:
        markdown_content: Markdown内容
        markdown_dir: Markdown文件所在目录
        
    Returns:
        图片路径列表
    """
    logger.info("[图片提取] 开始从Markdown内容中提取图片路径")
    logger.info(f"[图片提取] Markdown目录: {markdown_dir}")
    logger.info(f"[图片提取] Markdown内容长度: {len(markdown_content)}")
    
    # 匹配Markdown图片语法: ![alt](path)
    image_pattern = r'!\[.*?\]\((.*?)\)'
    matches = re.findall(image_pattern, markdown_content)
    
    logger.info(f"[图片提取] 找到 {len(matches)} 个图片引用")
    for i, match in enumerate(matches):
        logger.info(f"[图片提取] 图片引用 {i+1}: {match}")
    
    image_paths = []
    for match in matches:
        # 处理相对路径
        if match.startswith('images/'):
            # 相对于markdown文件的images目录
            image_path = os.path.join(markdown_dir, match)
        elif match.startswith('./images/'):
            # 相对于markdown文件的images目录
            image_path = os.path.join(markdown_dir, match[2:])  # 去掉 ./
        elif not os.path.isabs(match):
            # 其他相对路径
            image_path = os.path.join(markdown_dir, match)
        else:
            # 绝对路径
            image_path = match
            
        logger.info(f"[图片提取] 处理图片路径: {match} -> {image_path}")
        
        # 检查文件是否存在
        if os.path.exists(image_path):
            image_paths.append(image_path)
            logger.info(f"[图片提取] 图片文件存在: {image_path}")
        else:
            logger.warning(f"[图片提取] 图片文件不存在: {image_path}")
    
    logger.info(f"[图片提取] 最终提取到 {len(image_paths)} 个有效图片路径")
    for i, path in enumerate(image_paths):
        logger.info(f"[图片提取] 有效图片 {i+1}: {path}")
        
    return image_paths


def _create_image_mapping(image_paths: List[str]) -> Dict:
    """
    创建图片映射数据结构
    
    Args:
        image_paths: 图片路径列表
        
    Returns:
        图片映射字典
    """
    mapping = {"slide_1": {"slide_number": 1, "images": []}}
    
    for i, image_path in enumerate(image_paths):
        filename = os.path.basename(image_path)
        image_info = {
            "filename": filename,
            "filepath": image_path,
            "shape_id": i,
            "left": 0,
            "top": 0,
            "width": 0,
            "height": 0,
            "content_type": "",
            "original_format": "",
            "file_size": os.path.getsize(image_path) if os.path.exists(image_path) else 0
        }
        mapping["slide_1"]["images"].append(image_info)
    
    return mapping


def _translate_text(text: str, target_language: str, source_language: str) -> str:
    """
    翻译文本
    
    Args:
        text: 待翻译文本
        target_language: 目标语言
        source_language: 源语言
        
    Returns:
        翻译后的文本，如果翻译失败返回None
    """
    try:
        # 使用QwenTranslator进行翻译
        translator = QwenTranslator(target_language=target_language)
        translated = translator.translate_text(text, source_language=source_language)
        # 只有当翻译结果与原文不同时才返回翻译结果，否则返回None表示翻译失败
        if translated and translated.strip() and translated.strip() != text.strip():
            return translated
        else:
            return None
    except Exception as e:
        logger.error(f"翻译文本时出错: {str(e)}")
        logger.exception("翻译错误详情")
        return None


def process_markdown_images_ocr_and_translate(
    markdown_content: str,
    markdown_dir: str,
    target_language: str = "zh",
    source_language: str = "en"
) -> List[Dict]:
    """
    处理Markdown中的图片OCR识别和翻译
    
    Args:
        markdown_content: Markdown内容
        markdown_dir: Markdown文件所在目录
        target_language: 目标语言
        source_language: 源语言
        
    Returns:
        List of OCR结果字典，每个字典包含:
        {
            "success": bool,
            "image_path": str,
            "ocr_text_combined": str,
            "translation_text_combined": str
        }
    """
    logger.info("[Markdown OCR] 开始处理Markdown中的图片OCR识别和翻译")
    logger.info(f"[Markdown OCR] Markdown目录: {markdown_dir}")
    logger.info(f"[Markdown OCR] 目标语言: {target_language}, 源标语言: {source_language}")
    logger.info(f"[Markdown OCR] Markdown内容长度: {len(markdown_content)} 字符")
    
    results = []
    
    try:
        # 创建临时目录用于处理
        with tempfile.TemporaryDirectory(prefix="md_ocr_") as temp_dir:
            logger.info(f"[Markdown OCR] 创建临时目录: {temp_dir}")
            
            # 提取Markdown中的图片链接
            image_paths = _extract_image_paths_from_markdown(markdown_content, markdown_dir)
            logger.info(f"[Markdown OCR] 找到 {len(image_paths)} 个图片路径")
            
            if not image_paths:
                logger.info("[Markdown OCR] 未找到任何图片，跳过OCR处理")
                return results
            
            # 复制图片到临时目录
            temp_image_paths = []
            for i, image_path in enumerate(image_paths):
                if os.path.exists(image_path):
                    # 生成新的文件名
                    ext = os.path.splitext(image_path)[1]
                    new_filename = f"image_{i+1:04d}{ext}"
                    temp_image_path = os.path.join(temp_dir, new_filename)
                    
                    # 复制文件
                    shutil.copy2(image_path, temp_image_path)
                    temp_image_paths.append(temp_image_path)
                    logger.info(f"[Markdown OCR] 已复制图片: {os.path.basename(image_path)} -> {new_filename}")
                else:
                    logger.warning(f"[Markdown OCR] 图片文件不存在: {image_path}")
            
            if not temp_image_paths:
                logger.warning("[Markdown OCR] 没有有效的图片文件，跳过OCR处理")
                return results
            
            logger.info(f"[Markdown OCR] 成功复制 {len(temp_image_paths)} 个图片文件到临时目录")
            
            # 创建图片映射JSON文件
            image_mapping = _create_image_mapping(temp_image_paths)
            json_path = os.path.join(temp_dir, "image_mapping.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(image_mapping, f, ensure_ascii=False, indent=2)
            logger.info(f"[Markdown OCR] 已创建图片映射文件: {json_path}")
            logger.info(f"[Markdown OCR] 图片映射内容: {json.dumps(image_mapping, ensure_ascii=False, indent=2)}")
            
            # 执行OCR识别
            API_KEY = os.getenv("QWEN_API_KEY")
            if not API_KEY:
                logger.error("[Markdown OCR] 未设置QWEN_API_KEY环境变量，无法执行OCR识别")
                return results
            
            logger.info("[Markdown OCR] 开始执行OCR识别")
            process_folder_with_mapping(temp_dir, json_path, API_KEY)
            logger.info("[Markdown OCR] OCR识别完成")
            
            # 读取OCR结果
            if os.path.exists(json_path):
                with open(json_path, 'r', encoding='utf-8') as f:
                    updated_mapping = json.load(f)
                
                logger.info(f"[Markdown OCR] OCR结果文件加载成功")
                logger.info(f"[Markdown OCR] OCR结果内容: {json.dumps(updated_mapping, ensure_ascii=False, indent=2)}")
                
                # 处理每个图片的OCR结果
                for slide_key, slide_data in updated_mapping.items():
                    if 'images' not in slide_data:
                        logger.warning(f"[Markdown OCR] Slide数据中没有images字段: {slide_key}")
                        continue
                    
                    for image_info in slide_data['images']:
                        result = {
                            "image_path": image_info.get("filepath", ""),
                            "success": False,
                            "ocr_text_combined": "",
                            "translation_text_combined": ""
                        }
                        
                        # 获取OCR文本
                        all_text = image_info.get('all_text', {})
                        logger.info(f"[Markdown OCR] 图片 {os.path.basename(result['image_path'])} 的OCR文本: {all_text}")
                        
                        if all_text:
                            # 合并所有OCR文本
                            ocr_texts = [text.strip() for text in all_text.values() if text.strip()]
                            result["ocr_text_combined"] = '\n'.join(ocr_texts)
                            result["success"] = True
                            logger.info(f"[Markdown OCR] 图片OCR识别成功: {os.path.basename(result['image_path'])}, 文本长度: {len(result['ocr_text_combined'])}")
                        else:
                            logger.info(f"[Markdown OCR] 图片未识别到文本: {os.path.basename(result['image_path'])}")
                        
                        # 翻译OCR文本
                        if all_text:
                            try:
                                # 翻译所有文本
                                translated_texts = {}
                                for text_key, text_value in all_text.items():
                                    if text_value and text_value.strip():
                                        logger.info(f"[Markdown OCR] 翻译文本: {text_key} = {text_value.strip()}")
                                        translated = _translate_text(text_value.strip(), target_language, source_language)
                                        # 检查翻译是否成功
                                        if translated is not None:
                                            translated_texts[text_key] = translated
                                            logger.info(f"[Markdown OCR] 翻译结果: {text_key} = {translated}")
                                        else:
                                            logger.warning(f"[Markdown OCR] 翻译失败或结果与原文相同，跳过该文本段: {text_key}")
                                
                                # 合并所有翻译文本（包括未翻译的原文）
                                translation_texts = []
                                for text_key, text_value in all_text.items():
                                    if text_key in translated_texts:
                                        translation_texts.append(translated_texts[text_key].strip())
                                    elif text_value and text_value.strip():
                                        # 如果没有翻译成功，则使用原文
                                        translation_texts.append(text_value.strip())
                                
                                if translation_texts:
                                    result["translation_text_combined"] = '\n'.join(translation_texts)
                                    logger.info(f"[Markdown OCR] 图片翻译成功: {os.path.basename(result['image_path'])}, 翻译文本长度: {len(result['translation_text_combined'])}")
                                
                            except Exception as e:
                                logger.error(f"[Markdown OCR] 翻译图片文本时出错: {e}")
                                logger.exception("[Markdown OCR] 翻译错误详情")
                        
                        results.append(result)
                
                logger.info(f"[Markdown OCR] 处理完成，共处理 {len(results)} 个图片")
                for i, result in enumerate(results):
                    logger.info(f"[Markdown OCR] 结果 {i+1}: 成功={result['success']}, "
                              f"图片={os.path.basename(result['image_path'])}, "
                              f"OCR文本长度={len(result['ocr_text_combined'])}, "
                              f"翻译文本长度={len(result['translation_text_combined'])}")
            else:
                logger.error(f"[Markdown OCR] OCR结果文件不存在: {json_path}")
            
            return results
            
    except Exception as e:
        logger.error(f"[Markdown OCR] 处理Markdown图片OCR和翻译时出错: {str(e)}")
        logger.exception("[Markdown OCR] 详细错误信息")
        return results


class TextLineSplitter:
    """JSON文本行分割器 - 将包含换行符的文本拆分成多个text字段"""
    
    def __init__(self):
        self.processed_count = 0
        self.split_count = 0
    
    def process_json_file(self, json_file_path: str) -> bool:
        """
        处理JSON文件中的文本行分割
        
        Args:
            json_file_path: JSON文件路径
            
        Returns:
            是否处理成功
        """
        try:
            # 检查文件是否存在
            if not os.path.exists(json_file_path):
                logger.error(f"JSON文件不存在: {json_file_path}")
                return False
            
            # 读取JSON文件
            with open(json_file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # 处理数据
            self._process_mapping_data(data)
            
            # 保存处理后的数据
            with open(json_file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            logger.info(f"✅ 文本行分割完成！处理图片数: {self.processed_count}, 分割文本数: {self.split_count}")
            
            return True
            
        except Exception as e:
            logger.error(f"处理JSON文件时出错: {str(e)}")
            return False
    
    def _process_mapping_data(self, data: Dict[str, Any]):
        """处理映射数据中的所有文本"""
        for slide_key, slide_data in data.items():
            if 'images' not in slide_data:
                continue
            
            for image_info in slide_data['images']:
                filename = image_info.get('filename', '未知文件')
                
                # 处理all_text字段
                if 'all_text' in image_info and image_info['all_text']:
                    original_all_text = image_info['all_text'].copy()
                    new_all_text = self._split_text_dict(original_all_text)
                    
                    if new_all_text != original_all_text:
                        image_info['all_text'] = new_all_text
                        self.split_count += 1
                        logger.info(f"🔄 已分割 {filename} 的文本行")
                
                self.processed_count += 1
    
    def _split_text_dict(self, text_dict: Dict[str, str]) -> Dict[str, str]:
        """
        分割文本字典中的文本行
        
        Args:
            text_dict: 包含文本的字典
            
        Returns:
            分割后的文本字典
        """
        new_dict = {}
        text_counter = 1
        
        for key, text_value in text_dict.items():
            if not text_value:
                # 空文本直接跳过
                continue
            
            # 检查是否包含换行符
            if '\n' in text_value:
                # 按换行符分割
                lines = text_value.split('\n')
                
                logger.info(f"📝 发现多行文本 {key}: {len(lines)} 行")
                
                # 为每一行创建新的text字段
                for line in lines:
                    new_key = f"text{text_counter}"
                    new_dict[new_key] = line  # 保留空行，不进行strip()
                    text_counter += 1
                    
            else:
                # 单行文本，直接使用新的编号
                new_key = f"text{text_counter}"
                new_dict[new_key] = text_value
                text_counter += 1
        
        return new_dict


class PPTImageExtractor:
    """PPT图片提取器"""
    
    def __init__(self, temp_dir: str = None):
        self.temp_dir = temp_dir or tempfile.mkdtemp(prefix="ppt_ocr_")
        self.image_mapping = {}
        self.image_counter = 0
    
    def extract_images_from_slides(self, presentation_path: str, 
                                 selected_pages: Optional[List[int]] = None) -> Tuple[str, Dict]:
        """
        从PPT指定页面提取图片
        
        Args:
            presentation_path: PPT文件路径
            selected_pages: 选择的页面列表，None表示全选
            
        Returns:
            Tuple[临时文件夹路径, 图片映射字典]
        """
        try:
            prs = Presentation(presentation_path)
            total_slides = len(prs.slides)
            
            # 确定要处理的页面
            if selected_pages is None:
                pages_to_process = list(range(total_slides))
            else:
                pages_to_process = [p for p in selected_pages if 0 <= p < total_slides]
            
            logger.info(f"正在处理 {len(pages_to_process)} 个页面的图片...")
            
            for slide_idx in pages_to_process:
                slide = prs.slides[slide_idx]
                self._extract_images_from_slide(slide, slide_idx)
            
            # 保存映射关系到JSON文件
            mapping_file = os.path.join(self.temp_dir, "image_mapping.json")
            with open(mapping_file, 'w', encoding='utf-8') as f:
                json.dump(self.image_mapping, f, ensure_ascii=False, indent=2)
            
            logger.info(f"提取完成，共提取 {self.image_counter} 张图片")
            logger.info(f"临时文件夹: {self.temp_dir}")
            
            return self.temp_dir, self.image_mapping
            
        except Exception as e:
            logger.error(f"提取图片时发生错误: {str(e)}")
            raise
    
    def _extract_images_from_slide(self, slide, slide_idx: int):
        """从单个幻灯片提取图片"""
        slide_images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_info = self._save_image_from_shape(shape, slide_idx)
                    if image_info:
                        slide_images.append(image_info)
                except Exception as e:
                    logger.warning(f"提取第{slide_idx+1}页图片时出错: {str(e)}")
        
        if slide_images:
            self.image_mapping[f"slide_{slide_idx}"] = {
                "slide_number": slide_idx + 1,
                "images": slide_images
            }
    
    def _save_image_from_shape(self, shape: Picture, slide_idx: int) -> Optional[Dict]:
        """从图片形状保存图片文件"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            if not image_bytes:
                logger.warning(f"第{slide_idx+1}页图片数据为空，跳过")
                return None
            
            # 检测图片的真实格式
            actual_format, content_type = self._detect_image_format(image_bytes)
            
            # 根据真实格式确定文件扩展名
            ext_map = {
                'image/jpeg': '.jpg',
                'image/png': '.png', 
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
                'image/webp': '.webp',
                'image/x-emf': '.emf',
                'image/x-wmf': '.wmf',
                'unknown': '.bin'  # 未知格式保存为二进制文件
            }
            
            ext = ext_map.get(content_type, '.bin')
            
            # 生成唯一的文件名
            self.image_counter += 1
            filename = f"image_{self.image_counter:04d}{ext}"
            filepath = os.path.join(self.temp_dir, filename)
            
            # 如果是特殊格式，尝试转换为PNG
            if content_type in ['image/x-emf', 'image/x-wmf', 'unknown']:
                converted_filepath = self._try_convert_to_png(image_bytes, filename, filepath)
                if converted_filepath:
                    filename = os.path.basename(converted_filepath)
                    filepath = converted_filepath
                    content_type = 'image/png'
                    ext = '.png'
                else:
                    # 保存原始图片文件
                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)
            else:
                # 保存原始图片文件
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
            
            # 验证保存的文件是否有效
            if not self._validate_image_file(filepath):
                logger.warning(f"警告: 图片文件可能损坏 - {filename}")
                # 仍然返回信息，但标记为可能有问题
            
            logger.info(f"已保存图片: {filename} (格式: {actual_format}, 大小: {len(image_bytes)} bytes)")
            
            # 获取图片在幻灯片中的位置和大小信息
            return {
                "filename": filename,
                "filepath": filepath,
                "shape_id": shape.shape_id,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "content_type": content_type,
                "original_format": actual_format,
                "file_size": len(image_bytes)
            }
            
        except Exception as e:
            logger.error(f"保存图片时出错: {str(e)}")
            return None

    def _detect_image_format(self, image_bytes: bytes) -> tuple:
        """检测图片的真实格式"""
        try:
            # 检查文件头魔术字节
            if len(image_bytes) < 8:
                return "unknown", "unknown"
            
            header = image_bytes[:8]
            
            # PNG: 89 50 4E 47 0D 0A 1A 0A
            if header.startswith(b'\x89PNG\r\n\x1a\n'):
                return "PNG", "image/png"
            
            # JPEG: FF D8 FF
            elif header.startswith(b'\xff\xd8\xff'):
                return "JPEG", "image/jpeg"
            
            # GIF: 47 49 46 38 (GIF8)
            elif header.startswith(b'GIF8'):
                return "GIF", "image/gif"
            
            # BMP: 42 4D (BM)
            elif header.startswith(b'BM'):
                return "BMP", "image/bmp"
            
            # TIFF: 4D 4D 或 49 49
            elif header.startswith(b'MM') or header.startswith(b'II'):
                return "TIFF", "image/tiff"
            
            # WebP: 52 49 46 46 ... 57 45 42 50
            elif header.startswith(b'RIFF') and b'WEBP' in image_bytes[:12]:
                return "WebP", "image/webp"
            
            # EMF: 01 00 00 00 ... (Enhanced Metafile)
            elif len(image_bytes) > 40 and image_bytes[40:44] == b' EMF':
                return "EMF", "image/x-emf"
            
            # WMF: D7 CD C6 9A 或其他WMF标识
            elif header.startswith(b'\xd7\xcd\xc6\x9a') or header.startswith(b'\x01\x00\x09\x00'):
                return "WMF", "image/x-wmf"
            
            else:
                return "unknown", "unknown"
                
        except Exception as e:
            logger.error(f"检测图片格式时出错: {str(e)}")
            return "unknown", "unknown"

    def _try_convert_to_png(self, image_bytes: bytes, original_filename: str, original_filepath: str) -> Optional[str]:
        """尝试将特殊格式转换为PNG"""
        try:
            # 尝试使用PIL转换
            try:
                from PIL import Image
                import io
                
                # 先保存原始文件
                with open(original_filepath, 'wb') as f:
                    f.write(image_bytes)
                
                # 尝试用PIL打开并转换
                img = Image.open(io.BytesIO(image_bytes))
                
                # 转换为PNG
                png_filename = original_filename.replace(os.path.splitext(original_filename)[1], '_converted.png')
                png_filepath = os.path.join(os.path.dirname(original_filepath), png_filename)
                
                img.save(png_filepath, 'PNG')
                
                logger.info(f"已转换特殊格式图片为PNG: {png_filename}")
                return png_filepath
                
            except ImportError:
                logger.warning("警告: 未安装PIL/Pillow，无法转换特殊格式图片")
                return None
            except Exception as e:
                logger.error(f"PIL转换失败: {str(e)}")
                return None
                
        except Exception as e:
            logger.error(f"转换图片格式时出错: {str(e)}")
            return None

    def _validate_image_file(self, filepath: str) -> bool:
        """验证图片文件是否有效"""
        try:
            if not os.path.exists(filepath):
                return False
            
            file_size = os.path.getsize(filepath)
            if file_size == 0:
                return False
            
            # 尝试读取文件头
            with open(filepath, 'rb') as f:
                header = f.read(16)
                if len(header) < 4:
                    return False
            
            # 如果安装了PIL，尝试打开图片验证
            try:
                from PIL import Image
                img = Image.open(filepath)
                img.verify()  # 验证图片完整性
                return True
            except ImportError:
                # 没有PIL，只进行基本检查
                return True
            except Exception:
                return False
                
        except Exception:
            return False
    
    def cleanup(self):
        """清理临时文件"""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
            logger.info(f"已清理临时文件夹: {self.temp_dir}")


class PPTImageReplacer:
    """PPT图片文本添加器"""
    
    @staticmethod
    def add_ocr_text_to_slides(presentation_path: str, 
                             image_mapping: Dict,
                             output_path: str = None,
                             show_translation: bool = True):
        """
        在PPT页面右侧添加OCR识别的文本和翻译，保留原图片
        
        Args:
            presentation_path: 原PPT路径
            image_mapping: 包含all_text和translated_text字段的图片映射关系
            output_path: 输出PPT路径，None则覆盖原文件
            show_translation: 是否显示翻译结果
        """
        try:
            prs = Presentation(presentation_path)
            
            # 遍历每个幻灯片
            for slide_key, slide_info in image_mapping.items():
                slide_idx = slide_info["slide_number"] - 1
                slide = prs.slides[slide_idx]
                
                # 收集该页所有图片的OCR文本和翻译（改进版）
                slide_ocr_data = []
                image_count = 0
                
                for image_info in slide_info["images"]:
                    image_count += 1
                    filename = image_info.get("filename", f"图片{image_count}")
                    
                    # 获取原文和翻译的字典
                    original_texts = image_info.get('all_text', {})
                    translated_texts = image_info.get('translated_text', {}) if show_translation else {}
                    
                    # 构建文本对列表（原文-译文配对）
                    text_pairs = []
                    
                    if original_texts:
                        for key, original_text in original_texts.items():
                            # 确保文本不是空或"0"
                            cleaned_text = original_text.strip() if original_text else ""
                            if cleaned_text and cleaned_text != "0":
                                pair = {
                                    'key': key,
                                    'original': cleaned_text,
                                    'translated': ''
                                }
                                
                                # 查找对应的翻译，确保翻译文本也不是空或"0"
                                if key in translated_texts:
                                    cleaned_translation = translated_texts[key].strip() if translated_texts[key] else ""
                                    if cleaned_translation and cleaned_translation != "0":
                                        pair['translated'] = cleaned_translation
                                
                                text_pairs.append(pair)
                    
                    # 如果有文本对，添加到OCR数据中
                    if text_pairs:
                        ocr_data = {
                            'filename': filename,
                            'text_pairs': text_pairs
                        }
                        slide_ocr_data.append(ocr_data)
                
                # 如果该页有OCR数据，则添加文本框
                if slide_ocr_data:
                    PPTImageReplacer._add_paired_text_box_to_slide(
                        slide, slide_ocr_data, slide_info["slide_number"], show_translation
                    )
            
            # 保存PPT
            save_path = output_path or presentation_path
            prs.save(save_path)
            logger.info(f"OCR结果和翻译已添加到PPT: {save_path}")
            
        except Exception as e:
            logger.error(f"添加OCR文本和翻译时发生错误: {str(e)}")
            raise
    
    @staticmethod
    def _add_paired_text_box_to_slide(slide, ocr_data_list: List[Dict], slide_number: int, show_translation: bool):
        """在幻灯片右侧添加成对显示的OCR文本框"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        try:
            # 获取幻灯片尺寸
            try:
                presentation = slide.part.package.presentation_part.presentation
                slide_width = presentation.slide_width
                slide_height = presentation.slide_height
                
                # 转换为英寸
                slide_width_inches = slide_width / 914400
                slide_height_inches = slide_height / 914400
                
                logger.info(f"✅ 成功获取幻灯片尺寸: {slide_width_inches:.2f} x {slide_height_inches:.2f} 英寸")
                
            except Exception as size_error:
                logger.warning(f"⚠️ 无法获取幻灯片尺寸，使用标准尺寸: {str(size_error)}")
                # 使用标准的16:9幻灯片尺寸作为备用
                slide_width_inches = 13.33  # 标准宽屏PPT宽度
                slide_height_inches = 7.5   # 标准宽屏PPT高度
            
            # 计算文本框位置和大小
            if show_translation:
                # 有翻译时，文本框需要更宽
                text_box_width = Inches(slide_width_inches * 0.45)
            else:
                # 只有原文时，文本框稍窄
                text_box_width = Inches(slide_width_inches * 0.35)
            
            # 文本框高度：约为幻灯片高度的85%
            text_box_height = Inches(slide_height_inches * 0.85)
            
            # 文本框位置：在幻灯片右侧外部
            text_box_left = Inches(slide_width_inches + 0.1)
            text_box_top = Inches(0.3)
            
            # 添加文本框
            textbox = slide.shapes.add_textbox(
                text_box_left, 
                text_box_top, 
                text_box_width, 
                text_box_height
            )
            
            # 设置文本内容
            text_frame = textbox.text_frame
            text_frame.clear()  # 清除默认段落
            
            # 添加标题段落
            title_text = f"第 {slide_number} 页 OCR识别结果"
            if show_translation:
                title_text += " (逐行对照)"
            
            title_paragraph = text_frame.paragraphs[0]
            title_paragraph.text = title_text
            title_paragraph.font.size = Pt(14)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 黑色
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # 添加分隔线段落
            separator_paragraph = text_frame.add_paragraph()
            separator_paragraph.text = "=" * 40
            separator_paragraph.font.size = Pt(10)
            separator_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            separator_paragraph.alignment = PP_ALIGN.CENTER
            
            # 添加OCR内容段落（新的逐行对照显示逻辑）
            for i, ocr_data in enumerate(ocr_data_list, 1):
                filename = ocr_data['filename']
                text_pairs = ocr_data['text_pairs']
                
                # 图片标题段落
                img_title_paragraph = text_frame.add_paragraph()
                img_title_paragraph.text = f"\n🖼️ 图片 {i}: {filename}"
                img_title_paragraph.font.size = Pt(12)
                img_title_paragraph.font.bold = True
                img_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                # 遍历每个文本对，交错显示原文和译文
                for j, text_pair in enumerate(text_pairs, 1):
                    original_text = text_pair['original']
                    translated_text = text_pair['translated']
                    
                    # 原文部分
                    if original_text:
                        # 原文标签
                        original_label_paragraph = text_frame.add_paragraph()
                        original_label_paragraph.text = "🔤 原文:"
                        original_label_paragraph.font.size = Pt(9)
                        original_label_paragraph.font.bold = True
                        original_label_paragraph.font.color.rgb = RGBColor(0, 102, 204)  # 蓝色
                        
                        # 原文内容
                        original_content_paragraph = text_frame.add_paragraph()
                        original_content_paragraph.text = original_text
                        original_content_paragraph.font.size = Pt(9)
                        original_content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # 翻译部分（紧跟在对应原文后面）
                    if show_translation and translated_text:
                        # 翻译标签
                        translated_label_paragraph = text_frame.add_paragraph()
                        translated_label_paragraph.text = "🌐 译文:"
                        translated_label_paragraph.font.size = Pt(9)
                        translated_label_paragraph.font.bold = True
                        translated_label_paragraph.font.color.rgb = RGBColor(204, 102, 0)  # 橙色
                        
                        # 翻译内容
                        translated_content_paragraph = text_frame.add_paragraph()
                        translated_content_paragraph.text = translated_text
                        translated_content_paragraph.font.size = Pt(9)
                        translated_content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # 如果不是最后一个文本对，添加小分隔线
                    if j < len(text_pairs):
                        small_sep_paragraph = text_frame.add_paragraph()
                        small_sep_paragraph.text = "- - -"
                        small_sep_paragraph.font.size = Pt(7)
                        small_sep_paragraph.font.color.rgb = RGBColor(180, 180, 180)  # 更浅的灰色
                        small_sep_paragraph.alignment = PP_ALIGN.CENTER
                
                # 图片之间的分隔线
                if i < len(ocr_data_list):  # 不是最后一个图片则添加分隔线
                    sep_paragraph = text_frame.add_paragraph()
                    sep_paragraph.text = "─" * 35
                    sep_paragraph.font.size = Pt(10)
                    sep_paragraph.font.color.rgb = RGBColor(128, 128, 128)  # 灰色分隔线
                    sep_paragraph.alignment = PP_ALIGN.CENTER
            
            # 设置文本框边框和背景
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(248, 248, 248)  # 浅灰色背景
            
            # 设置边框
            textbox.line.color.rgb = RGBColor(200, 200, 200)  # 浅灰色边框
            textbox.line.width = Pt(1)
            
            # 设置文本框边距
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # 设置自动换行
            text_frame.word_wrap = True
            
            # 统计显示的文本对数量
            total_pairs = sum(len(ocr_data['text_pairs']) for ocr_data in ocr_data_list)
            translation_info = "和翻译" if show_translation else ""
            logger.info(f"✅ 已在第{slide_number}页右侧添加OCR文本框{translation_info}")
            logger.info(f"   📊 包含{len(ocr_data_list)}张图片，共{total_pairs}个文本对")
            
        except Exception as e:
            logger.error(f"在第{slide_number}页添加成对文本框时出错: {str(e)}")
            # 备用简单方案
            try:
                # 使用固定位置作为备用方案
                textbox = slide.shapes.add_textbox(
                    Inches(11),  # 固定位置：右侧
                    Inches(1),   # 固定位置：顶部
                    Inches(4),   # 固定宽度
                    Inches(6)    # 固定高度
                )
                
                text_frame = textbox.text_frame
                simple_content = f"第{slide_number}页OCR结果 (逐行对照):\n\n"
                
                for i, ocr_data in enumerate(ocr_data_list, 1):
                    simple_content += f"{i}. {ocr_data['filename']}\n"
                    for j, text_pair in enumerate(ocr_data['text_pairs'], 1):
                        simple_content += f"  {j}. 原文: {text_pair['original']}\n"
                        if show_translation and text_pair['translated']:
                            simple_content += f"     译文: {text_pair['translated']}\n"
                        simple_content += "\n"
                
                text_frame.text = simple_content
                text_frame.paragraphs[0].font.size = Pt(9)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                logger.warning(f"⚠️ 使用备用方案在第{slide_number}页添加OCR文本")
                
            except Exception as backup_error:
                logger.error(f"❌ 备用方案也失败: {str(backup_error)}")


def ocr_controller(presentation_path: str, 
                  selected_pages: Optional[List[int]] = None, 
                  output_path: str = None,
                  enable_translation: bool = True,
                  target_language: str = "中文",
                  source_language: str = "英文",
                  enable_text_splitting: str = "False") -> str:
    """
    OCR主控制器：提取图片、OCR识别、文本行分割、翻译、写回PPT
    
    Args:
        presentation_path: PPT文件路径
        selected_pages: 真实页码（第一页为1），None表示全篇处理
        output_path: 输出文件路径
        enable_translation: 是否启用翻译功能
        target_language: 目标语言
        source_language: 源语言
        enable_text_splitting: 是否启用文本行分割处理
        
    Returns:
        处理后的PPT文件路径
    """
    extractor = None
    ocr_processor = None
    try:
        # 验证输入文件
        if not os.path.exists(presentation_path):
            raise FileNotFoundError(f"PPT文件不存在: {presentation_path}")

        # 修正selected_pages为0-based索引
        prs = Presentation(presentation_path)
        total_slides = len(prs.slides)
        if selected_pages is not None:
            selected_pages = [p-1 for p in selected_pages if 1 <= p <= total_slides]
            if not selected_pages:
                logger.warning("selected_pages参数无有效页码，将处理全部页面")
                selected_pages = None

        # 1. 提取图片
        logger.info("=" * 50)
        logger.info("🔍 第一步：提取PPT中的图片")
        logger.info("=" * 50)
        extractor = PPTImageExtractor()
        temp_dir, image_mapping = extractor.extract_images_from_slides(
            presentation_path, selected_pages
        )
        if not image_mapping:
            logger.warning("未找到需要处理的图片")
            return presentation_path
        logger.info(f"✅ 图片提取完成，临时目录: {temp_dir}")

        # 2. 调用qwen-vl-ocr的api进行图片的文字提取
        logger.info("\n" + "=" * 50)
        logger.info("🤖 第二步：调用OCR QWEN API进行文本识别")
        logger.info("=" * 50)

        folder_path = temp_dir  # 替换为你的图片文件夹路径
        json_path = os.path.join(temp_dir, "image_mapping.json")  # 使用temp_dir作为文件夹路径
        API_KEY = os.getenv("QWEN_API_KEY")

        # 执行批量处理
        process_folder_with_mapping(folder_path, json_path, API_KEY)

        # 3. 文本行分割处理（可选）
        if enable_text_splitting == "True_spliting":
            logger.info("\n" + "=" * 50)
            logger.info("✂️ 第三步：文本行分割处理")
            logger.info("=" * 50)
            
            splitter = TextLineSplitter()
            split_success = splitter.process_json_file(json_path)
            
            if split_success:
                logger.info("✅ 文本行分割完成")
            else:
                logger.warning("⚠️ 文本行分割失败，将使用原始文本继续处理")
        else:
            logger.info("\n" + "=" * 50)
            logger.info("⏭️ 第三步：跳过文本行分割处理")
            logger.info("=" * 50)
            logger.info("✅ 保持原始文本格式")

        # 4. 翻译OCR识别结果
        step_num = 4 if enable_text_splitting != "False" else 3
        if enable_translation:
            logger.info("\n" + "=" * 50)
            logger.info(f"🌐 第{step_num}步：翻译识别结果 ({source_language} → {target_language})")
            logger.info("=" * 50)
            
            translation_success = TranslationManager.translate_ocr_results(
                temp_dir=temp_dir,
                target_language=target_language,
                source_language=source_language
            )
            
            if translation_success:
                logger.info(f"✅ 翻译完成")
                
                # 显示翻译摘要
                mapping_file = os.path.join(temp_dir, "image_mapping.json")
                summary = TranslationManager.get_translation_summary(mapping_file)
                if summary:
                    logger.info(f"📊 翻译摘要:")
                    logger.info(f"   - 总图片数: {summary.get('total_images', 0)}")
                    logger.info(f"   - 包含文本的图片: {summary.get('images_with_text', 0)}")
                    logger.info(f"   - 包含翻译的图片: {summary.get('images_with_translation', 0)}")
                    logger.info(f"   - 翻译成功率: {summary.get('translation_success_rate', 0):.1f}%")
            else:
                logger.warning("⚠️ 翻译失败，将只显示原文")
                enable_translation = False

        # 5. 读取更新后的映射文件
        step_num = 5 if enable_translation else (4 if enable_text_splitting != "False" else 3)
        logger.info(f"\n" + "=" * 50)
        logger.info(f"📖 第{step_num}步：读取处理结果")
        logger.info("=" * 50)
        mapping_file = os.path.join(temp_dir, "image_mapping.json")
        if not os.path.exists(mapping_file):
            raise Exception(f"映射文件不存在: {mapping_file}")
        with open(mapping_file, 'r', encoding='utf-8') as f:
            updated_mapping = json.load(f)
        logger.info("✅ 处理结果读取完成")
        
        # 统计结果
        ocr_count = 0
        translation_count = 0
        for slide_info in updated_mapping.values():
            for image_info in slide_info.get("images", []):
                if "all_text" in image_info and image_info["all_text"]:
                    ocr_count += 1
                    filename = image_info.get("filename", "未知文件")
                    text_preview = str(list(image_info["all_text"].values())[0])[:50] + "..."
                    logger.info(f"   📄 {filename}: {text_preview}")
                    
                    if enable_translation and "translated_text" in image_info and image_info["translated_text"]:
                        translation_count += 1
                        trans_preview = str(list(image_info["translated_text"].values())[0])[:50] + "..."
                        logger.info(f"   🌐 翻译: {trans_preview}")
        
        logger.info(f"📊 共识别出 {ocr_count} 张包含文本的图片")
        if enable_translation:
            logger.info(f"📊 共翻译了 {translation_count} 张图片的文本")

        # 6. 将OCR结果和翻译添加到PPT右侧
        step_num = 6 if enable_translation else (5 if enable_text_splitting != "False" else 4)
        logger.info(f"\n" + "=" * 50)
        content_desc = "OCR识别结果和翻译" if enable_translation else "OCR识别结果"
        logger.info(f"🎨 第{step_num}步：在PPT右侧添加{content_desc}")
        logger.info("=" * 50)
        
        PPTImageReplacer.add_ocr_text_to_slides(
            presentation_path=presentation_path,
            image_mapping=updated_mapping,
            output_path=output_path,
            show_translation=enable_translation
        )
        
        success_desc = "OCR结果和翻译" if enable_translation else "OCR结果"
        logger.info(f"✅ {success_desc}已添加到PPT右侧")
        logger.info("\n" + "=" * 50)
        logger.info("🎉 处理完成！")
        logger.info("=" * 50)
        return output_path or presentation_path
        
    except Exception as e:
        error_msg = f"OCR控制器处理失败: {str(e)}"
        logger.error(f"❌ {error_msg}")
        return presentation_path
    finally:
        if extractor:
            logger.info("🧹 清理临时文件...")
            extractor.cleanup()


# 使用示例
if __name__ == "__main__":
    # OCR API Token (替换为你的实际token)
    OCR_TOKEN = "你的OCR_API_TOKEN"
    QWEN_API_KEY = "你的QWEN_API_KEY"
    
    # 设置环境变量
    os.environ["OCR_TOKEN"] = OCR_TOKEN
    os.environ["QWEN_API_KEY"] = QWEN_API_KEY
    
    # PPT文件路径
    ppt_path = "example.pptx"
    
    # 选择要处理的页面 (1-based索引，None表示全部页面)
    selected_pages = [1, 2, 3]  # 处理前3页
    # selected_pages = None  # 处理全部页面
    
    # 输出文件路径
    output_path = "ocr_translated_result.pptx"
    
    # 执行OCR和翻译处理
    success = ocr_controller(
        presentation_path=ppt_path,
        selected_pages=selected_pages,
        output_path=output_path,
        enable_translation=True,  # 启用翻译
        target_language="英文",   # 目标语言
        source_language="中文",   # 源语言
        enable_text_splitting=True  # 启用文本行分割
    )
    
    if success:
        logger.info("🎉 PPT OCR和翻译处理成功！")
        logger.info(f"🔗 结果文件: {output_path}")
    else:
        logger.error("❌ 处理失败，请检查错误信息。")
