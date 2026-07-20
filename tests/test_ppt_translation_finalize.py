from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation


def test_finalize_translated_presentation_moves_translated_file_over_original(
    tmp_path: Path,
) -> None:
    # Given
    from app.function.ppt_translation_finalize import (
        FinalizePresentationRequest,
        finalize_translated_presentation,
    )

    original_path = tmp_path / "deck.pptx"
    translated_path = tmp_path / "deck_translated.pptx"
    original_path.write_text("original", encoding="utf-8")
    translated_path.write_text("translated", encoding="utf-8")
    request = FinalizePresentationRequest(
        translated_path=translated_path,
        original_path=original_path,
        selected_pages=(),
        source_language="English",
        target_language="Chinese",
        enable_text_splitting="False",
        progress_callback=None,
    )

    # When
    result = finalize_translated_presentation(request)

    # Then
    assert result is True
    assert original_path.read_text(encoding="utf-8") == "translated"
    assert not translated_path.exists()


def test_finalize_rejects_out_of_range_ocr_pages_instead_of_succeeding(
    tmp_path: Path,
) -> None:
    # Given
    from app.function.image_ocr.ocr_controller import InvalidSelectedPages
    from app.function.ppt_translation_finalize import (
        FinalizePresentationRequest,
        finalize_translated_presentation,
    )

    original_path = tmp_path / "deck.pptx"
    translated_path = tmp_path / "deck_translated.pptx"
    original_path.write_text("original", encoding="utf-8")
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(translated_path)
    request = FinalizePresentationRequest(
        translated_path=translated_path,
        original_path=original_path,
        selected_pages=(2,),
        source_language="English",
        target_language="Chinese",
        enable_text_splitting="True",
        progress_callback=None,
    )

    # When / Then
    with pytest.raises(InvalidSelectedPages):
        finalize_translated_presentation(request)

    assert original_path.read_text(encoding="utf-8") == "original"
    assert translated_path.exists()


def test_process_presentation_returns_after_xml_translation_before_layout(
    tmp_path: Path,
    monkeypatch,
) -> None:
    # Given
    from app.function import ppt_translate_async
    from app.function.pynuo_fuc import pyuno_controller as controller_module

    original_path = tmp_path / "deck.pptx"
    translated_path = tmp_path / "deck_xml_translated.pptx"
    original_path.write_text("original", encoding="utf-8")

    def fake_controller(*args, **kwargs):
        translated_path.write_text("translated", encoding="utf-8")
        return str(translated_path)

    async def fail_layout(_presentation_path: str) -> bool:
        raise AssertionError("layout adjustment should not run after XML translation succeeds")

    monkeypatch.setattr(controller_module, "pyuno_controller", fake_controller)
    monkeypatch.setattr(ppt_translate_async, "_adjust_ppt_layout_async", fail_layout)

    # When
    result = ppt_translate_async.process_presentation(
        str(original_path),
        stop_words_list=[],
        custom_translations={},
        select_page=[1],
        source_language="English",
        target_language="Chinese",
        bilingual_translation="paragraph_up",
        model="qwen",
        enable_text_splitting="False",
        enable_uno_conversion=True,
    )

    # Then
    assert result is True
    assert original_path.read_text(encoding="utf-8") == "translated"
