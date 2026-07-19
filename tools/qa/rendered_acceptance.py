from __future__ import annotations

import hashlib
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path

import fitz
import numpy as np
import psutil
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from app.function.pynuo_fuc.pptx_xml_translate import XmlTranslationRequest, translate_pptx_with_xml
from app.translation.libreoffice import LibreOfficeProcessAdapter, LibreOfficeRequest


@dataclass(frozen=True, slots=True)
class Box:
    left: float
    top: float
    width: float
    height: float

    def expanded(self, amount: float) -> fitz.Rect:
        return fitz.Rect(
            self.left - amount,
            self.top - amount,
            self.left + self.width + amount,
            self.top + self.height + amount,
        )


TARGET_BOX = Box(72, 72, 324, 144)
NEIGHBOR_BOX = Box(432, 72, 180, 144)


def create_overflow_fixture(path: Path) -> str:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    target = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4.5), Inches(2))
    target.text_frame.word_wrap = True
    target.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    target.text_frame.paragraphs[0].add_run().text = "ORIGINAL TARGET"
    target.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    neighbor = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(2.5), Inches(2))
    neighbor.text_frame.paragraphs[0].add_run().text = "PROTECTED NEIGHBOR"
    neighbor.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "UNCHANGED SECOND SLIDE"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    tokens = " ".join(f"TX{index:03d}" for index in range(45))
    return f"HEAD {tokens} TAIL"


def translate_first_unit(source: Path, output: Path, translated_text: str) -> Path:
    def deterministic(text_boxes, progress, source_language, target_language, model, stop_words, glossary):
        first = text_boxes[0]
        page = first["page_index"]
        key = f"{first['box_index'] + 1}_{first['paragraph_index'] + 1}"
        return {page: {"translated_fragments": {key: [translated_text]}}}

    request = XmlTranslationRequest(
        input_path=source,
        output_path=output,
        selected_page_indices=(0,),
        source_language="English",
        target_language="Chinese",
        model="qwen",
        stop_words=(),
        custom_translations={},
        bilingual_translation="translation_only",
        progress_callback=None,
    )
    return Path(translate_pptx_with_xml(request, translator=deterministic))


def render_presentation(
    presentation: Path,
    output_dir: Path,
    executable: Path,
    profile_root: Path,
) -> tuple[Path, int]:
    adapter = LibreOfficeProcessAdapter(executable, profile_root)
    result = adapter.convert(LibreOfficeRequest(presentation, output_dir))
    return result.output_path, result.pid


def audit_synthetic(source_pdf: Path, translated_pdf: Path, expected_text: str) -> dict[str, bool | float | int]:
    with fitz.open(source_pdf) as baseline, fitz.open(translated_pdf) as translated:
        page = translated[0]
        text = page.get_text()
        words = [word for word in page.get_text("words") if word[4] in {"HEAD", "TAIL"} or word[4].startswith("TX")]
        target = TARGET_BOX.expanded(2)
        neighbor = NEIGHBOR_BOX.expanded(0)
        inside = all(target.contains(fitz.Rect(word[:4])) for word in words)
        neighbor_overlap = sum(_intersection_area(fitz.Rect(word[:4]), neighbor) for word in words)
        sizes = _translated_font_sizes(page)
        baseline_pixels = _pixels(baseline[0])
        translated_pixels = _pixels(page)
        outside_ratio = _outside_diff_ratio(baseline_pixels, translated_pixels, TARGET_BOX)
        nonblank = float(np.any(translated_pixels < 248, axis=2).mean())
        return {
            "slide_count_preserved": len(baseline) == len(translated) == 2,
            "stable_page_dimensions": baseline[0].rect == translated[0].rect,
            "nonblank_pixel_ratio_positive": nonblank > 0.001,
            "all_markers_present_once": text.count("HEAD") == 1 and text.count("TAIL") == 1,
            "all_expected_tokens_present": all(text.count(token) == 1 for token in expected_text.split()),
            "all_glyphs_inside_textboxes": bool(words) and inside,
            "new_neighbor_overlap_sq_pt_lte_0_25": neighbor_overlap <= 0.25,
            "effective_font_pt_gte_8": bool(sizes) and min(sizes) >= 8.0,
            "outside_mask_diff_ratio_lte_0_0005": outside_ratio <= 0.0005,
            "nonblank_pixel_ratio": round(nonblank, 6),
            "neighbor_overlap_sq_pt": round(neighbor_overlap, 6),
            "effective_font_pt": round(min(sizes), 3) if sizes else 0.0,
            "outside_mask_diff_ratio": round(outside_ratio, 8),
        }


def slide_xml_equal(first: Path, second: Path, slide_number: int) -> bool:
    with zipfile.ZipFile(first) as left, zipfile.ZipFile(second) as right:
        name = f"ppt/slides/slide{slide_number}.xml"
        return left.read(name) == right.read(name)


def count_normal_autofit(path: Path, slide_number: int = 1) -> int:
    with zipfile.ZipFile(path) as archive:
        return archive.read(f"ppt/slides/slide{slide_number}.xml").count(b"normAutofit")


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def copy_source(source: Path, destination: Path) -> Path:
    destination.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, destination)
    return destination


def processes_exited(process_ids: list[int]) -> bool:
    return all(not psutil.pid_exists(process_id) for process_id in process_ids)


def _translated_font_sizes(page: fitz.Page) -> list[float]:
    sizes: list[float] = []
    for block in page.get_text("rawdict").get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = "".join(char.get("c", "") for char in span.get("chars", []))
                if "HEAD" in text or "TAIL" in text or "TX" in text:
                    sizes.append(float(span.get("size", 0)))
    return sizes


def _pixels(page: fitz.Page) -> np.ndarray:
    pixmap = page.get_pixmap(alpha=False)
    return np.frombuffer(pixmap.samples, dtype=np.uint8).reshape(pixmap.height, pixmap.width, pixmap.n)


def _outside_diff_ratio(first: np.ndarray, second: np.ndarray, box: Box) -> float:
    changed = np.any(np.abs(first.astype(np.int16) - second.astype(np.int16)) > 8, axis=2)
    mask = np.zeros(changed.shape, dtype=bool)
    left = max(0, int(box.left) - 2)
    top = max(0, int(box.top) - 2)
    right = min(mask.shape[1], int(box.left + box.width) + 2)
    bottom = min(mask.shape[0], int(box.top + box.height) + 2)
    mask[top:bottom, left:right] = True
    return float(np.logical_and(changed, ~mask).mean())


def _intersection_area(first: fitz.Rect, second: fitz.Rect) -> float:
    intersection = first & second
    return max(0.0, intersection.width) * max(0.0, intersection.height)
